"""Gera o deck de apresentação do módulo de opções (`opcoes/`).

Saída: ``apresentacao_modulo_opcoes.pptx`` na raiz do repositório.

Storytelling: Problema -> Solução -> Diferenciais. Público misto (técnico + negócio).
Todo o conteúdo (números e fórmulas) é fiel ao código-fonte do módulo.

Uso:
    python3 scripts/gerar_deck_opcoes.py
"""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# --- Paleta / tema ---------------------------------------------------------
AZUL_ESCURO = RGBColor(0x0B, 0x1F, 0x3A)   # fundo de capa / títulos
AZUL = RGBColor(0x1E, 0x5A, 0x9C)          # destaque primário
CIANO = RGBColor(0x2E, 0xC4, 0xB6)         # destaque / setas
VERDE = RGBColor(0x2E, 0x8B, 0x57)         # positivo
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)      # dor / risco
CINZA = RGBColor(0x4A, 0x4A, 0x4A)         # corpo de texto
CINZA_CLARO = RGBColor(0xEC, 0xF0, 0xF4)   # fundo de caixas
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

FONTE = "Calibri"
FONTE_MONO = "Consolas"

# Slide 16:9
LARGURA = Emu(12192000)
ALTURA = Emu(6858000)

ROOT = Path(__file__).resolve().parents[1]
ASSETS = Path(__file__).resolve().parent / "assets"


def _set_fill(shape, cor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()


def _add_textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def _run(par, texto, *, tam=18, cor=CINZA, negrito=False, italico=False, fonte=FONTE):
    r = par.add_run()
    r.text = texto
    r.font.size = Pt(tam)
    r.font.color.rgb = cor
    r.font.bold = negrito
    r.font.italic = italico
    r.font.name = fonte
    return r


def _faixa_topo(slide, titulo, kicker=None):
    """Faixa de título superior comum aos slides de conteúdo."""
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, Emu(150000))
    _set_fill(barra, CIANO)

    _, tf = _add_textbox(slide, Emu(500000), Emu(280000), Emu(11200000), Emu(900000))
    if kicker:
        p = tf.paragraphs[0]
        _run(p, kicker.upper(), tam=13, cor=AZUL, negrito=True)
    p = tf.add_paragraph() if kicker else tf.paragraphs[0]
    _run(p, titulo, tam=32, cor=AZUL_ESCURO, negrito=True)


def _fundo(slide, cor):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, ALTURA)
    _set_fill(bg, cor)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --- Slides especializados -------------------------------------------------
def slide_capa(prs):
    s = _blank(prs)
    _fundo(s, AZUL_ESCURO)

    # acento
    acento = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Emu(2600000), Emu(220000), Emu(1500000))
    _set_fill(acento, CIANO)

    _, tf = _add_textbox(s, Emu(700000), Emu(2400000), Emu(10500000), Emu(2200000))
    p = tf.paragraphs[0]
    _run(p, "MÓDULO DE OPÇÕES", tam=20, cor=CIANO, negrito=True)
    p = tf.add_paragraph()
    _run(p, "Screening inteligente de Covered Calls", tam=44, cor=BRANCO, negrito=True)
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    _run(
        p,
        "Da cadeia bruta de opções ao ranking acionável — probabilidade conservadora, "
        "custos reais e backtest, num pipeline vetorizado.",
        tam=18,
        cor=RGBColor(0xC8, 0xD4, 0xE0),
    )

    _, tf = _add_textbox(s, Emu(700000), Emu(5900000), Emu(10500000), Emu(500000))
    p = tf.paragraphs[0]
    _run(p, "plataforma allocation  ·  módulo opcoes  ·  2026-06-09", tam=14, cor=RGBColor(0x8A, 0x9B, 0xB0))


def slide_problema(prs):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, "Escolher uma covered call é uma decisão multidimensional", kicker="O problema")

    _, tf = _add_textbox(s, Emu(500000), Emu(1500000), Emu(11200000), Emu(700000))
    p = tf.paragraphs[0]
    _run(
        p,
        "Cada contrato combina strike, prazo, liquidez, volatilidade e custos. "
        "Nenhuma métrica isolada captura o trade-off — e o que parece bom numa dimensão é ruim em outra.",
        tam=18,
        cor=CINZA,
    )

    # dois cartões comparativos
    dados = [
        ("Strike OTM +5%", "Baixa probabilidade de exercício", "...mas o retorno do prêmio é fraco", VERMELHO),
        ("Strike OTM +10%", "Retorno do prêmio atraente", "...mas o risco de atribuição é alto demais", AZUL),
    ]
    x = Emu(700000)
    for titulo, linha1, linha2, cor in dados:
        cartao = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(2700000), Emu(5100000), Emu(2400000))
        _set_fill(cartao, CINZA_CLARO)
        cartao.line.color.rgb = cor
        cartao.line.width = Pt(2)
        tf = cartao.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(300000)
        tf.margin_top = Emu(250000)
        p = tf.paragraphs[0]
        _run(p, titulo, tam=24, cor=cor, negrito=True)
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        _run(p, "✓ " + linha1, tam=18, cor=VERDE, negrito=True)
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        _run(p, "✗ " + linha2, tam=18, cor=VERMELHO)
        x = Emu(6400000)

    _, tf = _add_textbox(s, Emu(700000), Emu(5400000), Emu(10800000), Emu(700000))
    p = tf.paragraphs[0]
    _run(p, "→ É preciso uma régua única que balanceie retorno, risco e custo.", tam=20, cor=AZUL_ESCURO, negrito=True)


def slide_bullets(prs, kicker, titulo, itens, *, subtitulo=None):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, titulo, kicker=kicker)

    y = 1500000
    if subtitulo:
        _, tf = _add_textbox(s, Emu(500000), Emu(y), Emu(11200000), Emu(600000))
        p = tf.paragraphs[0]
        _run(p, subtitulo, tam=18, cor=CINZA, italico=True)
        y += 750000

    _, tf = _add_textbox(s, Emu(700000), Emu(y), Emu(10800000), Emu(ALTURA - Emu(y) - Emu(500000)))
    for i, item in enumerate(itens):
        if isinstance(item, tuple):
            cabeca, corpo = item
        else:
            cabeca, corpo = item, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, "●  ", tam=18, cor=CIANO, negrito=True)
        _run(p, cabeca, tam=20, cor=AZUL_ESCURO, negrito=True)
        if corpo:
            _run(p, "  —  " + corpo, tam=18, cor=CINZA)
    return s


def slide_jornada(prs):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, "A jornada de uma opção pela engine", kicker="Solução · visão geral")

    etapas = [
        ("Dado bruto", "cadeia de opções\n(yfinance / mock)", AZUL),
        ("Pipeline", "30+ métricas\nvetorizadas", AZUL),
        ("Probabilidade", "max(d2, empírica)\nconservadora", CIANO),
        ("Custos", "prêmio líquido\nrisk-adjusted", CIANO),
        ("Score", "retorno × (1−prob)\n× theta / vega", VERDE),
        ("Top-N", "ranking + matriz\nauditável", VERDE),
    ]
    n = len(etapas)
    margem = 500000
    gap = 180000
    largura_total = 11192000
    w = (largura_total - (n - 1) * gap) / n
    y = 2500000
    h = 1700000
    for i, (titulo, corpo, cor) in enumerate(etapas):
        x = int(margem + i * (w + gap))
        cx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(int(w)), Emu(h))
        _set_fill(cx, cor)
        tf = cx.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, titulo, tam=17, cor=BRANCO, negrito=True)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        for j, linha in enumerate(corpo.split("\n")):
            if j > 0:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
            _run(p, linha, tam=12, cor=RGBColor(0xE6, 0xF7, 0xF4))
        # seta
        if i < n - 1:
            sx = int(x + w + gap * 0.1)
            seta = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(sx), Emu(y + h // 2 - 70000), Emu(int(gap * 0.8)), Emu(140000))
            _set_fill(seta, RGBColor(0xB0, 0xBE, 0xC5))

    _, tf = _add_textbox(s, Emu(500000), Emu(4700000), Emu(11200000), Emu(900000))
    p = tf.paragraphs[0]
    _run(p, "Camadas desacopladas: ", tam=18, cor=AZUL_ESCURO, negrito=True)
    _run(p, "data/ (provedor) → models/ (matemática pura) → pipeline.py → calls.py. ", tam=18, cor=CINZA)
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    _run(p, "Sem imports para cima — os mesmos modelos servirão para puts, spreads e collar.", tam=18, cor=CINZA)


def slide_formula(prs, kicker, titulo, formula_linhas, termos, *, rodape=None):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, titulo, kicker=kicker)

    # caixa de fórmula
    caixa = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(700000), Emu(1600000), Emu(10800000), Emu(1300000))
    _set_fill(caixa, AZUL_ESCURO)
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(300000)
    for i, linha in enumerate(formula_linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _run(p, linha, tam=20, cor=CIANO, negrito=True, fonte=FONTE_MONO)

    _, tf = _add_textbox(s, Emu(700000), Emu(3200000), Emu(10800000), Emu(3200000))
    for i, (termo, desc) in enumerate(termos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        _run(p, termo, tam=18, cor=AZUL, negrito=True, fonte=FONTE_MONO)
        _run(p, "  —  " + desc, tam=17, cor=CINZA)

    if rodape:
        _, tf = _add_textbox(s, Emu(700000), Emu(6150000), Emu(10800000), Emu(500000))
        p = tf.paragraphs[0]
        _run(p, rodape, tam=16, cor=VERDE, negrito=True, italico=True)


def slide_arquitetura(prs):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, "Arquitetura em camadas — separação rígida, sem imports para cima", kicker="Diferenciais")

    def caixa(x, y, w, h, titulo, sub, cor, txt_cor=BRANCO, tam=16):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
        _set_fill(c, cor)
        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, titulo, tam=tam, cor=txt_cor, negrito=True)
        if sub:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            _run(p, sub, tam=11, cor=txt_cor)
        return c

    # topo: CLI -> Config -> Runner
    caixa(700000, 1500000, 3400000, 800000, "CLI  ·  cli.py", "argparse, subcomandos", AZUL_ESCURO)
    caixa(4400000, 1500000, 3400000, 800000, "Config  ·  config.py", "TOML validado (dataclass)", AZUL_ESCURO)
    caixa(8100000, 1500000, 3400000, 800000, "Runner  ·  runner.py", "loop por ativo", AZUL_ESCURO)

    # camada de serviço
    servicos = [
        ("data/", "ProvedorDados\nyfinance ↔ mock", AZUL),
        ("models/", "Black-Scholes,\nGreeks, empírica", CIANO),
        ("opcoes/pipeline.py", "30+ métricas\nvetorizadas", AZUL),
        ("opcoes/calls.py", "filtros, score,\ntop-N", VERDE),
        ("report.py + risco/", "Excel, payoff,\ncarteira", AZUL),
    ]
    n = len(servicos)
    margem = 700000
    gap = 150000
    largura_total = 10792000
    w = (largura_total - (n - 1) * gap) / n
    y = 2900000
    for i, (titulo, sub, cor) in enumerate(servicos):
        x = int(margem + i * (w + gap))
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(int(w)), Emu(1600000))
        _set_fill(c, cor)
        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, titulo, tam=14, cor=BRANCO, negrito=True, fonte=FONTE_MONO)
        for linha in sub.split("\n"):
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            _run(p, linha, tam=11, cor=RGBColor(0xEA, 0xF3, 0xF0))

    caixa(700000, 4800000, 10792000, 700000,
          "models/ é matemática pura e reutilizável → mesma base alimenta puts, spreads, collar e backtest",
          None, CINZA_CLARO, txt_cor=AZUL_ESCURO, tam=15)


def slide_roadmap(prs):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, "Roadmap — stubs prontos para expansão", kicker="Diferenciais · futuro")

    _, tf = _add_textbox(s, Emu(500000), Emu(1450000), Emu(11200000), Emu(550000))
    p = tf.paragraphs[0]
    _run(p, "A arquitetura já reserva os pontos de extensão. Cada stub tem docstring com o escopo previsto.", tam=17, cor=CINZA, italico=True)

    itens = [
        ("opcoes/puts.py", "Venda de puts cash-secured / naked — espelho da covered call, fecha a estratégia Wheel", CIANO),
        ("opcoes/spreads.py", "Bull/bear spread, iron condor, calendar — crédito líquido, max profit/loss, break-even", AZUL),
        ("opcoes/volatilidade.py", "IV rank/percentile, term structure, skew, prêmio de volatilidade (IV − histórica)", AZUL),
        ("opcoes/hedge.py", "Collar, put protetora, beta hedge — proteção integrada à posição coberta", VERDE),
    ]
    y = 2150000
    for titulo, corpo, cor in itens:
        barra = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(700000), Emu(y), Emu(150000), Emu(900000))
        _set_fill(barra, cor)
        _, tf = _add_textbox(s, Emu(1000000), Emu(y), Emu(10400000), Emu(900000), anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        _run(p, titulo, tam=19, cor=AZUL_ESCURO, negrito=True, fonte=FONTE_MONO)
        p = tf.add_paragraph()
        _run(p, corpo, tam=16, cor=CINZA)
        y += 1050000


def slide_fechamento(prs):
    s = _blank(prs)
    _fundo(s, AZUL_ESCURO)

    acento = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Emu(2200000), Emu(220000), Emu(1900000))
    _set_fill(acento, CIANO)

    _, tf = _add_textbox(s, Emu(700000), Emu(2100000), Emu(10800000), Emu(2600000))
    p = tf.paragraphs[0]
    _run(p, "Uma régua única para covered calls", tam=40, cor=BRANCO, negrito=True)
    p = tf.add_paragraph()
    p.space_before = Pt(16)
    _run(
        p,
        "Varre → ranqueia → valida. Probabilidade conservadora, custos reais e backtest, "
        "com auditabilidade total e arquitetura pronta para puts, spreads e hedge.",
        tam=20,
        cor=RGBColor(0xC8, 0xD4, 0xE0),
    )
    p = tf.add_paragraph()
    p.space_before = Pt(24)
    _run(p, "allocation --config config.toml", tam=20, cor=CIANO, negrito=True, fonte=FONTE_MONO)


# --- Gráfico de payoff (matplotlib) ---------------------------------------
def gerar_payoff_png() -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    destino = ASSETS / "payoff_covered_call.png"

    s0, strike, premio = 100.0, 110.0, 3.0
    precos = np.linspace(80, 135, 300)
    acoes = precos - s0
    call_vendida = -np.maximum(precos - strike, 0) + premio
    coberta = acoes + call_vendida

    fig, ax = plt.subplots(figsize=(7.6, 4.3), dpi=150)
    ax.axhline(0, color="#9aa6b2", lw=1)
    ax.plot(precos, acoes, "--", color="#b0bec5", lw=2, label="Ação (buy & hold)")
    ax.plot(precos, coberta, color="#1E5A9C", lw=3, label="Covered call")
    ax.axvline(strike, color="#2EC4B6", ls=":", lw=2)
    ax.fill_between(precos, coberta, 0, where=(coberta >= 0), color="#2E8B57", alpha=0.10)
    ax.annotate("teto de ganho", xy=(122, premio + (strike - s0)), xytext=(118, 16),
                color="#C0392B", fontsize=11, fontweight="bold")
    ax.annotate("strike", xy=(strike, -18), xytext=(strike + 0.5, -19), color="#2EC4B6", fontsize=11, fontweight="bold")
    ax.set_xlabel("Preço do ativo no vencimento (S_T)", fontsize=11)
    ax.set_ylabel("Resultado por ação ($)", fontsize=11)
    ax.set_title("Payoff de uma covered call", fontsize=13, fontweight="bold", color="#0B1F3A")
    ax.legend(loc="upper left", frameon=False, fontsize=10)
    ax.grid(alpha=0.15)
    fig.tight_layout()
    fig.savefig(destino, bbox_inches="tight")
    plt.close(fig)
    return destino


def slide_payoff(prs, png_path):
    s = _blank(prs)
    _fundo(s, BRANCO)
    _faixa_topo(s, "O instrumento: troca upside por renda imediata", kicker="O problema · contexto")

    s.shapes.add_picture(str(png_path), Emu(600000), Emu(1700000), height=Emu(4400000))

    _, tf = _add_textbox(s, Emu(7400000), Emu(1900000), Emu(4400000), Emu(4200000), anchor=MSO_ANCHOR.MIDDLE)
    pontos = [
        "Vender a call gera prêmio imediato — renda sobre a posição.",
        "Em troca, o ganho fica limitado acima do strike (call exercida).",
        "O dilema: quanto teto sacrificar por quanto prêmio?",
        "O módulo transforma esse trade-off num número comparável.",
    ]
    for i, t in enumerate(pontos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        _run(p, "●  ", tam=16, cor=CIANO, negrito=True)
        _run(p, t, tam=17, cor=CINZA)


# --- Montagem --------------------------------------------------------------
def construir():
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA

    # ATO 1 — Problema
    slide_capa(prs)
    slide_problema(prs)
    slide_payoff(prs, gerar_payoff_png())
    slide_bullets(
        prs,
        "O problema · por que é difícil",
        "Cinco lacunas que ferramentas comuns deixam abertas",
        [
            ("Informação fragmentada", "strike, prazo, liquidez, vol e custos vivem em planilhas separadas"),
            ("Probabilidade heterogênea", "o risk-neutral (mercado) diverge da frequência histórica real do ativo"),
            ("Custos ignorados", "corretagem de venda e exercício mudam qual opção é de fato atrativa"),
            ("Atribuição antecipada", "calls americanas com dividendo alto e delta ≥ 0,70 podem ser exercidas cedo"),
            ("Sem validação", "ranquear é fácil; provar que o critério funciona no histórico, não"),
        ],
    )

    # ATO 2 — Solução
    slide_jornada(prs)
    slide_bullets(
        prs,
        "Solução · pipeline",
        "Pipeline vetorizado — 30+ métricas num passo só",
        [
            ("preparar_calls_para_modelo()", "núcleo em opcoes/pipeline.py; adicionar métrica = editar aqui"),
            ("NumPy / pandas, sem loops", "~100+ strikes de 20 ativos processados em segundos"),
            ("Liquidez, distância, retorno", "volume, open interest, spread; retorno do prêmio anualizado"),
            ("IV usada + fonte", "IV implícita preferida; fallback para vol histórica (fonte_vol)"),
            ("Greeks completos", "delta, gamma, vega, theta, rho — alimentam score e flags de risco"),
        ],
        subtitulo="Uma única função vetorizada computa todas as métricas por opção.",
    )
    slide_formula(
        prs,
        "Solução · probabilidade",
        "Probabilidade de exercício: o pior caso entre dois mundos",
        [
            "prob_final = max( prob_d2 ,  prob_empírica )",
            "d2 = [ ln(S₀/K) + (r − q − ½·IV²)·T ] / ( IV·√T )",
        ],
        [
            ("prob_d2 = N(d2)", "Black-Scholes risk-neutral: probabilidade implícita no preço de mercado"),
            ("prob_empírica", "frequência histórica em janelas NÃO-sobrepostas (independência estatística)"),
            ("max(...)", "escolhe a estimativa mais pessimista — postura conservadora por design"),
        ],
        rodape="5 anos de histórico + janelas de 7–20 dias → dezenas de amostras independentes.",
    )
    slide_formula(
        prs,
        "Solução · custos",
        "Modelo de custos: prêmio líquido ajustado ao risco",
        [
            "prêmio_líquido = prêmio − custo_venda − custo_exercício × prob_final",
            "custo_exercício = max( 0,25% × strike × 100 ,  US$ 10 )",
        ],
        [
            ("custo_venda", "descontado sempre, por ação"),
            ("custo_exercício × prob_final", "ponderado pela probabilidade de a call ser exercida"),
            ("buy-write", "inclui também o custo de compra das ações na base"),
        ],
        rodape="Custos explícitos mudam o ranking — sobretudo em strikes de baixo valor nocional.",
    )
    slide_formula(
        prs,
        "Solução · score",
        "O score: uma régua única de retorno × segurança",
        [
            "score = retorno_se_exercido_anual × (1 − prob_final)",
            "        × (1 + peso_theta·θ_eff) / (1 + peso_vega·vega_risk)",
        ],
        [
            ("retorno_se_exercido_anual", "ganho anualizado se a call for exercida (já líquido de custos)"),
            ("(1 − prob_final)", "penaliza o risco de atribuição — prêmio só vale se for 'seguro'"),
            ("peso_theta · θ_eff", "bônus por time decay; θ_eff = theta normalizado pelo prêmio líquido"),
            ("peso_vega · vega_risk", "penaliza exposição à volatilidade; pesos = 0 por padrão"),
        ],
        rodape="Default (pesos = 0): score = retorno_se_exercido × (1 − prob_final).",
    )
    slide_bullets(
        prs,
        "Solução · filtros + auditoria",
        "Cadeia de filtros ordenada — e nada se perde",
        [
            ("1 · Liquidez", "volume, open interest e spread dentro dos limites"),
            ("2 · Cap de probabilidade", "prob_exercicio_final ≤ prob_exerc_max"),
            ("3 · Distância de strike", "OTM mínimo configurável (min_distancia_strike_pct)"),
            ("4 · Retorno líquido > 0", "depois de todos os custos"),
            ("5 · Lucro se exercido > 0", "nunca recomenda vender no prejuízo"),
            ("Coluna status + matriz_opcoes.xlsx", "rejeitadas ficam registradas com o motivo — auditável"),
        ],
        subtitulo="Precedência clara: o primeiro motivo de reprovação 'vence' e é registrado.",
    )

    # ATO 3 — Diferenciais
    slide_bullets(
        prs,
        "Diferenciais · validação",
        "Backtest: o critério é provado no histórico",
        [
            ("backtest_covered_call()", "simula vendas sucessivas de covered calls ao longo do histórico"),
            ("Prêmio por Black-Scholes", "estimado com a volatilidade realizada de cada data de entrada"),
            ("Métricas de resultado", "retorno médio/anualizado, taxa de exercício, taxa de acerto, volatilidade"),
            ("Comparação com buy & hold", "mostra o que a estratégia agrega (ou não) sobre só carregar a ação"),
            ("Baseado em modelo", "valida a lógica de seleção de strike e horizonte — honesto quanto ao escopo"),
        ],
    )
    slide_bullets(
        prs,
        "Diferenciais · engenharia",
        "Por que é robusto e reutilizável",
        [
            ("Vetorização", "NumPy/pandas ponta a ponta — rápido e sem loops frágeis"),
            ("Fallback inteligente", "vol histórica quando falta IV; d2 quando a empírica é insuficiente"),
            ("Modo offline", "base_mock/ com 23 tickers — demos e testes sem rede"),
            ("Config declarativa (TOML)", "parâmetros e overrides por ativo sem tocar no código"),
            ("Logging estruturado", "rejeições rastreadas por motivo; erro num ativo não derruba os demais"),
        ],
    )
    slide_arquitetura(prs)
    slide_roadmap(prs)
    slide_fechamento(prs)

    destino = ROOT / "apresentacao_modulo_opcoes.pptx"
    prs.save(destino)
    print(f"Deck gerado: {destino}  ({len(prs.slides._sldIdLst)} slides)")
    return destino


if __name__ == "__main__":
    construir()
