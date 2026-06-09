"""Gera o deck de apresentação do Modelo de Alocação (All&In)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------- paleta
AZUL_ESCURO = RGBColor(0x0B, 0x1F, 0x3A)   # fundo títulos
AZUL        = RGBColor(0x1E, 0x5A, 0xA8)   # destaque
CIANO       = RGBColor(0x17, 0xB8, 0xC4)   # acento
CINZA       = RGBColor(0x3C, 0x44, 0x52)   # corpo
CINZA_CLARO = RGBColor(0x6B, 0x72, 0x80)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
FUNDO_CLARO = RGBColor(0xF4, 0xF7, 0xFB)
VERDE       = RGBColor(0x1E, 0xA0, 0x6B)

LARG = Inches(13.333)
ALT  = Inches(7.5)

prs = Presentation()
prs.slide_width = LARG
prs.slide_height = ALT
BLANK = prs.slide_layouts[6]


def slide_base(cor_fundo=BRANCO):
    s = prs.slides.add_slide(BLANK)
    fundo = s.shapes.add_shape(1, 0, 0, LARG, ALT)
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = cor_fundo
    fundo.line.fill.background()
    fundo.shadow.inherit = False
    s.shapes._spTree.remove(fundo._element)
    s.shapes._spTree.insert(2, fundo._element)
    return s


def caixa(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def par(tf, texto, tamanho, cor, bold=False, primeiro=False, align=PP_ALIGN.LEFT,
        space_after=6, italic=False, nivel=0):
    p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = nivel
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(tamanho)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = cor
    r.font.name = "Calibri"
    return p


def barra_lateral(slide, cor=CIANO, w=Inches(0.18)):
    b = slide.shapes.add_shape(1, 0, 0, w, ALT)
    b.fill.solid()
    b.fill.fore_color.rgb = cor
    b.line.fill.background()
    b.shadow.inherit = False


def chip(slide, x, y, texto, cor=CIANO):
    w = Inches(0.4 + 0.13 * len(texto))
    c = slide.shapes.add_shape(5, x, y, w, Inches(0.42))
    c.fill.solid()
    c.fill.fore_color.rgb = cor
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = BRANCO
    r.font.name = "Calibri"
    return c


def rodape(slide, num):
    _, tf = caixa(slide, Inches(11.8), Inches(7.0), Inches(1.4), Inches(0.4))
    par(tf, f"Modelo de Alocação · {num}", 9, CINZA_CLARO, primeiro=True,
        align=PP_ALIGN.RIGHT)


def slide_conteudo(titulo, etapa_chip, bullets, num, sub=None):
    """bullets: lista de (texto, nivel, bold)."""
    s = slide_base(BRANCO)
    barra_lateral(s)
    if etapa_chip:
        chip(s, Inches(0.55), Inches(0.5), etapa_chip, AZUL)
    _, tf = caixa(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(1.1))
    par(tf, titulo, 32, AZUL_ESCURO, bold=True, primeiro=True)
    if sub:
        par(tf, sub, 15, CINZA_CLARO, italic=True, space_after=0)
    _, tf2 = caixa(s, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.6))
    primeiro = True
    for texto, nivel, bold in bullets:
        cor = AZUL_ESCURO if bold and nivel == 0 else CINZA
        tam = 20 if nivel == 0 else 16
        marcador = "" if bold and nivel == 0 else ("•  " if nivel == 0 else "–  ")
        p = par(tf2, marcador + texto, tam, cor, bold=bold, primeiro=primeiro,
                nivel=nivel, space_after=10 if nivel == 0 else 5)
        primeiro = False
    rodape(s, num)
    return s


# ================================================================ SLIDE 1 — Capa
s = slide_base(AZUL_ESCURO)
faixa = s.shapes.add_shape(1, 0, Inches(4.6), LARG, Inches(0.08))
faixa.fill.solid(); faixa.fill.fore_color.rgb = CIANO; faixa.line.fill.background()
faixa.shadow.inherit = False
_, tf = caixa(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.5))
par(tf, "Do palpite à precisão", 50, BRANCO, bold=True, primeiro=True, space_after=4)
par(tf, "Como recomendamos a carteira certa para cada cliente — na escala de todos",
    24, CIANO, space_after=0)
_, tf = caixa(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.5))
par(tf, "Modelo de Alocação  ·  Avenue Intelligence", 20, BRANCO, bold=True, primeiro=True)
par(tf, "De uma planilha por cliente a um motor que escala para toda a base.",
    15, CINZA_CLARO, italic=True)

# ================================================================ SLIDE 2 — A dor
slide_conteudo(
    "O problema", "ATO 1 · PROBLEMA",
    [("Cada cliente é único", 0, True),
     ("Perfil de risco, restrições, necessidade de liquidez, objetivo de renda", 1, False),
     ("Posição atual e caixa disponível diferentes para cada um", 1, False),
     ("Recomendar \"na mão\" não escala", 0, True),
     ("Lento, inconsistente e difícil de auditar", 1, False),
     ("Conhecimento preso na cabeça de poucas pessoas", 1, False),
     ("A carteira deve estar em conformidade regulatória", 0, True),
     ("FINRA (EUA) e CVM (Brasil) impõem regras de suitability e adequação de produtos", 1, False),
     ("Qualquer recomendação precisa ser rastreável e justificável perante os reguladores", 1, False)],
    "2/14",
    sub="A pergunta que move este projeto")
# pergunta-gancho destacada
s = prs.slides[-1]
box = s.shapes.add_shape(5, Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.0))
box.fill.solid(); box.fill.fore_color.rgb = FUNDO_CLARO; box.line.color.rgb = CIANO
box.line.width = Pt(1.5); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "“Como dar a melhor recomendação para milhares de clientes ao mesmo tempo — cada um com suas regras?”"
r.font.size = Pt(18); r.font.bold = True; r.font.italic = True
r.font.color.rgb = AZUL; r.font.name = "Calibri"

# ================================================================ SLIDE 3 — Expectativa
slide_conteudo(
    "O que o cliente espera", "ATO 1 · PROBLEMA",
    [("Uma carteira aderente ao seu perfil (model portfolio)", 0, True),
     ("Respeito às suas restrições e à política da casa", 0, True),
     ("Os produtos de melhor qualidade primeiro", 0, True),
     ("Uma ordem de compra pronta para o caixa que ele tem hoje", 0, True),
     ("Facilidade de execução", 0, True),
     ("Receber exatamente o que comprar, quanto e em qual produto — sem ambiguidade", 1, False),
     ("Acompanhamento contínuo", 0, True),
     ("Saber se a carteira ainda está aderente ao alvo após movimentos de mercado", 1, False),
     ("Ser avisado quando houver oportunidade de melhora ou desvio relevante", 1, False)],
    "3/14")

# ================================================================ SLIDE 4 — Funil
s = slide_base(BRANCO)
barra_lateral(s)
chip(s, Inches(0.55), Inches(0.5), "ATO 2 · SOLUÇÃO", CIANO)
_, tf = caixa(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(1.0))
par(tf, "Visão geral — um funil de 5 estágios", 32, AZUL_ESCURO, bold=True, primeiro=True)
etapas = [
    ("1 · Produtos\nRecomendados", "O QUE pode\nser recomendado"),
    ("2 · Model\nPortfolio", "QUANTO de\ncada categoria"),
    ("3 · Restrições", "O que NÃO\npode entrar"),
    ("4 · Produtos\npor Cliente", "Lista sob\nmedida"),
    ("5 · Otimizador", "Ordem de\ncompra"),
]
x = Inches(0.6)
w = Inches(2.25)
gap = Inches(0.12)
y = Inches(2.7)
cores = [AZUL, AZUL, AZUL, AZUL, CIANO]
for i, (titulo, desc) in enumerate(etapas):
    card = s.shapes.add_shape(5, x, y, w, Inches(2.2))
    card.fill.solid(); card.fill.fore_color.rgb = cores[i]; card.line.fill.background()
    card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = titulo
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = BRANCO; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0xD6, 0xE6, 0xF7); r2.font.name = "Calibri"
    if i < 4:
        seta = s.shapes.add_shape(13, x + w + Emu(1), y + Inches(0.9),
                                  gap, Inches(0.4))
        seta.fill.solid(); seta.fill.fore_color.rgb = CINZA_CLARO; seta.line.fill.background()
        seta.shadow.inherit = False
    x = x + w + gap
_, tf = caixa(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.0))
par(tf, "Cada etapa estreita o universo até sobrar exatamente o que faz sentido para aquele cliente.",
    18, CINZA, italic=True, primeiro=True, align=PP_ALIGN.CENTER)
rodape(s, "4/14")

# ================================================================ SLIDE 5 — Etapa 1
slide_conteudo(
    "Etapa 1 — O que pode ser recomendado", "PRODUTOS RECOMENDADOS",
    [("Research seleciona, por mês, os melhores produtos", 0, True),
     ("Bonds, ETFs, Fundos e UCITS (ações não entram no modelo)", 1, False),
     ("Campanhas pontuais entram sem bagunçar a fila", 0, True),
     ("Recomendações de campanha \"furam a fila\" na posição certa", 1, False),
     ("Os produtos mensais são repriorizados automaticamente", 1, False),
     ("Cada produto chega com seus atributos", 0, True),
     ("Gestora, duration, setor, emissor, rating, aplicação mínima", 1, False)],
    "5/14")

# ================================================================ SLIDE 6 — Etapa 2
slide_conteudo(
    "Etapa 2 — Quanto de cada coisa", "MODEL PORTFOLIO",
    [("Peso-alvo por categoria, definido por perfil de risco", 0, True),
     ("Três carteiras disponíveis", 0, True),
     ("Agregada  ·  Expandida  ·  Personalizada", 1, False),
     ("Versionado no tempo", 0, True),
     ("Sabemos qual era o alvo recomendado em qualquer data", 1, False),
     ("A Personalizada lê a composição direto do IPS do cliente", 1, False)],
    "6/14")

# ================================================================ SLIDE 7 — Etapa 3
slide_conteudo(
    "Etapa 3 — As regras do jogo", "RESTRIÇÕES",
    [("Três camadas de proteção", 0, True),
     ("Produtos — o que o cliente não quer (via IPS)", 1, False),
     ("Atributos — duration, gestora, setor, emissor; liquidez e renda por política", 1, False),
     ("Concentração — limites de exposição por rating, duração e emissor", 1, False),
     ("Origem: preenchimento do IPS + política da casa", 0, True)],
    "7/14")

# ================================================================ SLIDE 8 — Etapa 4
slide_conteudo(
    "Etapa 4 — A lista sob medida", "PRODUTOS POR CLIENTE",
    [("Cruza o alvo (model portfolio) com os produtos recomendados", 0, True),
     ("Aplica as restrições em cadeia", 0, True),
     ("R1 produtos  →  R2 atributos  →  R3 concentração", 1, False),
     ("Ordena por prioridade final", 0, True),
     ("Qualidade do produto (ordem) + ranking da recomendação mensal", 1, False),
     ("No fim, cada cliente tem sua própria prateleira, já ordenada", 0, True),
     ("Cada versão diária é guardada para histórico e auditoria", 1, False)],
    "8/14")

# ================================================================ SLIDE 9 — Etapa 5 (herói)
s = slide_base(AZUL_ESCURO)
chip(s, Inches(0.55), Inches(0.5), "OTIMIZADOR", CIANO)
_, tf = caixa(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.8))
par(tf, "Etapa 5 — A decisão: water-filling", 30, BRANCO, bold=True, primeiro=True, space_after=2)
par(tf, "O caixa vai primeiro para quem está mais longe do alvo", 14, CIANO, italic=True, space_after=0)

# ---- explicação (lado esquerdo) ----
_, tf = caixa(s, Inches(0.55), Inches(2.0), Inches(4.2), Inches(4.8))
itens = [
    ("Water-filling", "Fecha os maiores gaps primeiro — sem nunca vender"),
    ("Respeita limites", "Aplicação mínima + teto de concentração sobre posição já existente"),
    ("Regra de sobra", "Resíduo vai para AV1010 (caixa/RF curto), com força se necessário"),
]
primeiro = True
for t, d in itens:
    par(tf, t, 16, BRANCO, bold=True, primeiro=primeiro, space_after=1)
    par(tf, d, 12, RGBColor(0xC9, 0xD6, 0xE6), space_after=12)
    primeiro = False

# ---- tabela antes / depois ----
# dados: Categoria | Atual% | Alvo% | Gap | Compra | Depois%
linhas = [
    ("Categoria",       "Atual", "Alvo",  "Gap",    "Compra",    "Depois"),
    ("AV1010  RF Curto","22%",   "15%",   "–",      "US$ 0",     "20%"),
    ("AV2020  RF Longo","8%",    "20%",   "12%",    "US$ 36k",   "20%"),
    ("AV3030  Ações",   "10%",   "15%",   "5%",     "US$ 15k",   "15%"),
    ("AV4040  Fundos",  "18%",   "25%",   "7%",     "US$ 21k",   "25%"),
    ("AV5050  Alternat.","4%",   "10%",   "6%",     "US$ 18k",   "10%"),
    ("Caixa",           "US$300k","–",   "–",       "–US$90k",   "US$210k"),
]
col_w = [Inches(2.2), Inches(0.7), Inches(0.7), Inches(0.65), Inches(0.95), Inches(0.85)]
row_h = Inches(0.46)
tx = Inches(4.85)
ty = Inches(2.0)
header_cor = AZUL
row_cores  = [RGBColor(0x14, 0x2A, 0x50), RGBColor(0x0F, 0x24, 0x45)]
for ri, row in enumerate(linhas):
    cx = tx
    for ci, cell in enumerate(row):
        is_header = ri == 0
        bg = header_cor if is_header else row_cores[ri % 2]
        rect = s.shapes.add_shape(1, cx, ty, col_w[ci], row_h)
        rect.fill.solid(); rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(0x1E, 0x3A, 0x6E); rect.line.width = Pt(0.5)
        rect.shadow.inherit = False
        tf = rect.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = cell
        r.font.size = Pt(11 if is_header else 12)
        r.font.bold = is_header
        # destaca coluna "Depois" em verde
        if ci == 5 and not is_header:
            r.font.color.rgb = RGBColor(0x4C, 0xE8, 0xA8)
        elif ci == 4 and not is_header and cell != "US$ 0" and cell != "–":
            r.font.color.rgb = CIANO
        else:
            r.font.color.rgb = BRANCO
        r.font.name = "Calibri"
        cx += col_w[ci]
    ty += row_h

# legenda
_, tf = caixa(s, Inches(4.85), Inches(5.3), Inches(8.4), Inches(0.4))
par(tf, "Exemplo ilustrativo — cliente perfil moderado, aporte de US$ 100k + caixa existente de US$ 200k",
    10, CINZA_CLARO, italic=True, primeiro=True)
rodape(s, "9/14")

# ================================================================ SLIDE 10 — Exemplo
slide_conteudo(
    "Exemplo real — antes e depois", "ATO 3 · IMPACTO",
    [("Cliente com aporte de US$ 100 mil, perfil moderado", 0, True),
     ("O modelo calcula peso atual vs. recomendado vs. depois do aporte", 1, False),
     ("Entrega uma ordem de compra pronta", 0, True),
     ("Produto a produto, com taxas, YTW e peso final no portfólio", 1, False),
     ("Caixa alocado nas categorias mais distantes do alvo, respeitando os limites", 1, False)],
    "10/14",
    sub="Da carteira atual à recomendação executável")

# ================================================================ SLIDE 11 — Por que importa
slide_conteudo(
    "Por que isso importa", "ATO 3 · IMPACTO",
    [("Escala — toda a base, todo dia, sem esforço manual", 0, True),
     ("Consistência — a mesma regra para todos, auditável", 0, True),
     ("Rastreabilidade — snapshots históricos de cada decisão", 0, True),
     ("Governança — política e research codificados, não na cabeça de alguém", 0, True)],
    "11/14")

# ================================================================ SLIDE 12 — Engenharia
slide_conteudo(
    "A engenharia que sustenta", "ATO 3 · IMPACTO",
    [("100% versionado em Git (Databricks Repos)", 0, True),
     ("Mudanças revisadas em PR, reversíveis", 1, False),
     ("Otimizador em Python puro", 0, True),
     ("Portável e testável fora do Databricks", 1, False),
     ("Camadas desacopladas", 0, True),
     ("Evoluir uma etapa sem quebrar as outras", 1, False)],
    "12/14")

# ================================================================ SLIDE 13 — Próximos passos
slide_conteudo(
    "Próximos passos", "ROADMAP",
    [("Reativar Stocks quando fizer sentido para o modelo", 0, True),
     ("Corrigir a ingestão da aplicação mínima de fundos", 0, True),
     ("Evoluir o otimizador", 0, True),
     ("Rebalanceamento com venda e restrições adicionais", 1, False)],
    "13/14")

# ================================================================ SLIDE 14 — Encerramento
s = slide_base(AZUL_ESCURO)
faixa = s.shapes.add_shape(1, 0, Inches(4.3), LARG, Inches(0.08))
faixa.fill.solid(); faixa.fill.fore_color.rgb = CIANO; faixa.line.fill.background()
faixa.shadow.inherit = False
_, tf = caixa(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.8), anchor=MSO_ANCHOR.MIDDLE)
par(tf, "Cada cliente, a carteira certa —", 40, BRANCO, bold=True, primeiro=True, space_after=2)
par(tf, "na escala de todos.", 40, CIANO, bold=True, space_after=0)
_, tf = caixa(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.8))
par(tf, "Obrigado  ·  Perguntas?", 20, RGBColor(0xC9, 0xD6, 0xE6), primeiro=True)

# ---------------------------------------------------------------- salvar
out = "Modelo_de_Alocacao.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
